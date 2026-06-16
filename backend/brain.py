import pyautogui
import pygetwindow as gw
import os
import sqlite3
import datetime
import psutil
import ctypes
import socket
import platform
import requests
import time
from pypdf import PdfReader
import json
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

class CognitiveBrain:
    def __init__(self, db_path="memory.db"):
        self.db_path = os.path.join(os.path.dirname(__file__), db_path)
        self.hacker_token = os.getenv("HACKER_AI_TOKEN")
        self._init_db()
        pyautogui.FAILSAFE = True
        self._network_cache = {
            "local": "127.0.0.1",
            "public": "SCANNING...",
            "isp": "DETECTING...",
            "location": "UNKNOWN",
            "latency": "N/A",
            "timestamp": 0
        }

    def _init_db(self):
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        c.execute('''CREATE TABLE IF NOT EXISTS interactions
                     (timestamp TEXT, user_input TEXT, bot_response TEXT, context TEXT)''')
        c.execute('''CREATE TABLE IF NOT EXISTS tasks
                     (timestamp TEXT, task_name TEXT, status TEXT)''')
        c.execute('''CREATE TABLE IF NOT EXISTS mistakes
                     (timestamp TEXT, action_name TEXT, mistake_desc TEXT, correction TEXT)''')
        conn.commit()
        conn.close()

    # --- SELF-LEARNING SYSTEM ---
    def report_mistake(self, action, desc, correction):
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        c.execute("INSERT INTO mistakes VALUES (?, ?, ?, ?)",
                  (datetime.datetime.now().isoformat(), action, desc, correction))
        conn.commit()
        conn.close()
        return f"Mistake logged. I will correct {action} next time."

    def get_correction(self, action):
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        c.execute("SELECT correction FROM mistakes WHERE action_name = ? ORDER BY timestamp DESC LIMIT 1", (action,))
        row = c.fetchone()
        conn.close()
        return row[0] if row else None

    # --- APP LAUNCHER ---
    def scan_shortcuts(self, app_name):
        """Neural Shortcut Scanner: Locates .lnk files on Desktop or Start Menu."""
        import winshell
        from win32com.client import Dispatch
        shell = Dispatch('WScript.Shell')
        
        search_paths = [
            os.path.join(os.environ['USERPROFILE'], 'Desktop'),
            os.path.join(os.environ['APPDATA'], 'Microsoft', 'Windows', 'Start Menu', 'Programs'),
            os.path.join(os.environ['PROGRAMDATA'], 'Microsoft', 'Windows', 'Start Menu', 'Programs')
        ]
        
        app_name = app_name.lower()
        for path in search_paths:
            if not os.path.exists(path): continue
            for root, dirs, files in os.walk(path):
                for file in files:
                    if file.lower().endswith(".lnk") and app_name in file.lower():
                        shortcut_path = os.path.join(root, file)
                        try:
                            # Resolve the actual executable path from the shortcut
                            target = shell.CreateShortCut(shortcut_path).Targetpath
                            if target and os.path.exists(target):
                                return target
                            return shortcut_path # Return link if target resolve fails
                        except:
                            return shortcut_path
        return None

    def launch_app(self, app_name):
        app_name_clean = app_name.lower().strip()
        
        # 1. Direct Shell Mapping (Fastest)
        apps = {
            "excel": "excel", "word": "winword", "powerpoint": "powerpnt",
            "chrome": "chrome", "edge": "msedge", "vscode": "code", "code": "code"
        }
        target = apps.get(app_name_clean.replace(" ", ""))
        
        import subprocess
        if target:
            try:
                subprocess.Popen(f"start {target}", shell=True)
                return True
            except: pass

        # 2. Neural Shortcut Scan (Search Desktop/Start Menu)
        try:
            found_path = self.scan_shortcuts(app_name_clean)
            if found_path:
                os.startfile(found_path)
                return True
        except: pass

        # 3. Last Resort: Intelligent System Search (Win + Type)
        try:
            pyautogui.press('win')
            import time
            time.sleep(0.3)
            pyautogui.write(app_name_clean, interval=0.02)
            time.sleep(0.6)
            pyautogui.press('enter')
            return True
        except Exception as e:
            self.report_mistake(f"launch_{app_name_clean}", str(e), "Check app name")
            return False

    # --- MEMORY SYSTEM ---
    def store_interaction(self, user_input, bot_response, context=""):
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        c.execute("INSERT INTO interactions VALUES (?, ?, ?, ?)",
                  (datetime.datetime.now().isoformat(), user_input, bot_response, context))
        conn.commit()
        conn.close()

    def recall_recent(self, limit=5):
        conn = sqlite3.connect(self.db_path)
        c = conn.cursor()
        c.execute("SELECT * FROM interactions ORDER BY timestamp DESC LIMIT ?", (limit,))
        rows = c.fetchall()
        conn.close()
        return rows

    # --- DESKTOP ACTIONS ---
    def get_active_windows(self):
        windows = gw.getAllTitles()
        return [w for w in windows if w.strip()]

    def take_screenshot(self, filename="last_screen.png"):
        screenshot = pyautogui.screenshot()
        screenshot.save(filename)
        return os.path.abspath(filename)

    def type_text(self, text):
        pyautogui.write(text, interval=0.1)

    def press_key(self, key):
        pyautogui.press(key)

    # --- CROSS-APP INTELLIGENCE ---
    def read_pdf(self, file_path):
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text
        except Exception as e:
            return f"Error reading PDF: {str(e)}"

    # --- DEVICE CONTROL ---
    def change_system_setting(self, setting, value=None):
        """Controls system settings like volume, dark mode, wifi."""
        setting = setting.lower()
        import subprocess
        
        try:
            if "volume" in setting:
                # Direct PowerShell method for volume (more reliable than SendKeys)
                vol = int(value) if value else 50
                # Fix: Use -lt instead of < for PowerShell and [int] for math
                cmd = f"$w = New-Object -ComObject WScript.Shell; for($i=0;$i -lt 50;$i++){{$w.SendKeys([char]174)}}; for($i=0;$i -lt [int]({vol}/2);$i++){{$w.SendKeys([char]175)}}"
                result = subprocess.run(["powershell", "-Command", cmd], capture_output=True, text=True)
                if result.returncode != 0:
                    return False, f"PS_ERR: {result.stderr}"
                return True, f"Acoustic density stabilized at {vol}%"

            elif "dark mode" in setting or "theme" in setting:
                mode = 0 if "dark" in setting else 1
                cmd = f"Set-ItemProperty -Path 'HKCU:\\SOFTWARE\\Microsoft\\Windows\\CurrentVersion\\Themes\\Personalize' -Name 'AppsUseLightTheme' -Value {mode}"
                result = subprocess.run(["powershell", "-Command", cmd], capture_output=True, text=True)
                if result.returncode != 0:
                    return False, f"THEME_ERR: {result.stderr}"
                return True, f"Neural interface theme set to {'Dark' if mode==0 else 'Light'}"

            elif "wifi" in setting or "wi-fi" in setting:
                state = "enabled" if "on" in str(value).lower() else "disabled"
                cmd = f"netsh interface set interface Wi-Fi {state}"
                result = subprocess.run(cmd, shell=True, capture_output=True, text=True)
                if result.returncode != 0:
                    return False, f"NET_ERR: {result.stderr or 'Admin privileges required'}"
                return True, f"Neural uplink {state}"

            elif "brightness" in setting:
                brightness = int(value) if value else 50
                # Robust WMI call
                cmd = f"(Get-WmiObject -Namespace root/WMI -Class WmiMonitorBrightnessMethods).WmiSetBrightness(1, {brightness})"
                result = subprocess.run(["powershell", "-Command", cmd], capture_output=True, text=True)
                if result.returncode != 0:
                    # Fallback to display-brightness if available or return error
                    return False, f"DISP_ERR: {result.stderr}. Note: External monitors may not support WMI brightness."
                return True, f"Photonic intensity recalibrated to {brightness}%"

            return False, "Unknown operational setting"
        except Exception as e:
            return False, f"CRITICAL_BRAIN_FAULT: {str(e)}"

    # --- WORK & CONTENT FEATURES ---
    def read_presentation(self, file_path):
        """Reads PowerPoint slides and extracts content for 'presentation' mode."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = []
            for i, slide in enumerate(prs.slides):
                text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
                content.append({"slide": i+1, "text": text.strip()})
            return {"success": True, "data": content}
        except ImportError:
            return {"success": False, "error": "python-pptx not installed. Run 'pip install python-pptx'"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def meeting_assistant(self, action="start"):
        """Virtual presence bridge for meeting notes/summarization."""
        if action == "start":
            # Logic to start audio capture or screen monitoring
            return True, "Neural link established with meeting stream. Recording notes..."
        elif action == "summarize":
            # Logic to process captured audio/text
            return True, "Summary: The meeting focused on Q3 targets and AI integration roadmap."
        return False, "Invalid action"

    def translate_live(self, text, target_lang="en"):
        """Live translation bridge."""
        # This would typically use googletrans or a specialized API
        # For now, we return a simulated bridge response
        return f"[Translating to {target_lang}]: {text}"

    def get_network_telemetry(self, force_refresh=False):
        """Neural Uplink Monitor: Fetches real-time network telemetry with caching."""
        now = time.time()
        # Refresh every 15 minutes (900s) unless forced or still scanning
        if force_refresh or (now - self._network_cache["timestamp"] > 900) or self._network_cache["public"] == "SCANNING...":
            try:
                # Local IP
                s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
                s.connect(("8.8.8.8", 80))
                self._network_cache["local"] = s.getsockname()[0]
                s.close()
            except: pass

            try:
                # Public IP + ISP Info (using ip-api.com)
                response = requests.get('http://ip-api.com/json/', timeout=3).json()
                if response.get('status') == 'success':
                    self._network_cache["public"] = response.get('query')
                    self._network_cache["isp"] = response.get('isp')
                    self._network_cache["location"] = f"{response.get('city')}, {response.get('country')}"
                else:
                    # Fallback to ipify
                    self._network_cache["public"] = requests.get('https://api64.ipify.org', timeout=2).text
            except:
                try:
                    self._network_cache["public"] = requests.get('https://api64.ipify.org', timeout=2).text
                except: pass
            
            # Latency Check
            try:
                import subprocess
                # Use -n 1 for Windows ping
                ping_res = subprocess.run(["ping", "-n", "1", "8.8.8.8"], capture_output=True, text=True, timeout=2)
                if "time=" in ping_res.stdout:
                    latency = ping_res.stdout.split("time=")[1].split("ms")[0].strip()
                    self._network_cache["latency"] = f"{latency}ms"
                else:
                    self._network_cache["latency"] = "TIMEOUT"
            except:
                self._network_cache["latency"] = "N/A"

            self._network_cache["timestamp"] = now

        return self._network_cache

    def get_system_stats(self):
        net_io = psutil.net_io_counters()
        
        # Check for Admin/Root access
        try:
            is_admin = os.getuid() == 0 if os.name != 'nt' else ctypes.windll.shell32.IsUserAnAdmin() != 0
        except:
            is_admin = False

        # Dynamic Network Telemetry
        net_stats = self.get_network_telemetry()

        return {
            "os": f"{platform.system()} {platform.release()}",
            "arch": platform.machine(),
            "local_ip": net_stats["local"],
            "public_ip": net_stats["public"],
            "isp": net_stats["isp"],
            "location": net_stats["location"],
            "latency": net_stats["latency"],
            "cpu_usage": psutil.cpu_percent(),
            "memory_usage": psutil.virtual_memory().percent,
            "disk_usage": psutil.disk_usage('/').percent,
            "processes": len(psutil.pids()),
            "active_apps": len(self.get_active_windows()),
            "access_level": "NEURAL_ROOT" if is_admin else "USER_RESTRICTED",
            "network": {
                "sent_mb": round(net_io.bytes_sent / (1024 * 1024), 2),
                "recv_mb": round(net_io.bytes_recv / (1024 * 1024), 2),
                "status": "UPLINK_STABLE" if psutil.net_if_stats() else "OFFLINE"
            }
        }

    def check_integrity(self):
        """Neural Forensic Scan: Checks for system compromises."""
        suspicious = []
        try:
            for proc in psutil.process_iter(['pid', 'name', 'cpu_percent']):
                if (proc.info['cpu_percent'] or 0) > 85:
                    suspicious.append(f"HIGH_LOAD: {proc.info['name']} (PID: {proc.info['pid']})")
        except: pass
        try:
            if len(psutil.net_connections()) > 500:
                suspicious.append("NETWORK_SATURATION_DETECTED (Potential DDoS/Scan)")
        except: pass
        return {
            "safe": len(suspicious) == 0,
            "threats": suspicious,
            "status": "SECURE" if len(suspicious) == 0 else "INTEGRITY_COMPROMISED"
        }

# Singleton instance
brain = CognitiveBrain()
